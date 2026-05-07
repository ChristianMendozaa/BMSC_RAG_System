import logging
from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


@dataclass
class TextBlock:
    text: str
    page_number: int
    block_type: str = "text"


@dataclass
class ImageBlock:
    data: bytes
    page_number: int
    image_index: int


@dataclass
class ParseResult:
    text_blocks: list[TextBlock] = field(default_factory=list)
    image_blocks: list[ImageBlock] = field(default_factory=list)


def parse(file_bytes: bytes) -> ParseResult:
    result = ParseResult()
    image_idx = 0

    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        logger.error("Failed to open PPTX: %s", exc)
        raise

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            slide_texts: list[str] = []

            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                slide_texts.append(line)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_data = shape.image.blob
                            result.image_blocks.append(
                                ImageBlock(
                                    data=img_data,
                                    page_number=slide_num,
                                    image_index=image_idx,
                                )
                            )
                            image_idx += 1
                        except Exception as exc:
                            logger.warning(
                                "Skipping image on slide %d: %s", slide_num, exc
                            )
                except Exception as exc:
                    logger.warning(
                        "Skipping shape on slide %d: %s", slide_num, exc
                    )

            try:
                notes = slide.notes_slide
                notes_text = notes.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[Notas del presentador]: {notes_text}")
            except Exception:
                pass

            if slide_texts:
                result.text_blocks.append(
                    TextBlock(
                        text=f"[Slide {slide_num}]\n" + "\n".join(slide_texts),
                        page_number=slide_num,
                    )
                )
        except Exception as exc:
            logger.warning("Skipping PPTX slide %d: %s", slide_num, exc)

    return result
