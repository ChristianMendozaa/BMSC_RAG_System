"""Tests de los parsers de ingesta y del dispatch parse_file.

Las fixtures se generan programáticamente (no hay binarios commiteados).
No cargan modelos ni tocan DB/ChromaDB: solo importan app.services.parsers.
"""

from io import BytesIO

import fitz
import pytest
from docx import Document as DocxDocument
from docx.shared import Inches as DocxInches
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches

from app.services.parsers import (
    ACCEPTED_EXTENSIONS,
    docx_parser,
    image_parser,
    parse_file,
    pdf_parser,
    pptx_parser,
    xlsx_parser,
)

PNG_MAGIC = b"\x89PNG"


# ── Generadores de fixtures ───────────────────────────────────────────────


def make_png(width: int = 120, height: int = 120, color=(200, 30, 30)) -> bytes:
    buf = BytesIO()
    Image.new("RGB", (width, height), color).save(buf, format="PNG")
    return buf.getvalue()


def make_pdf() -> bytes:
    doc = fitz.open()
    page = doc.new_page()  # A4: 595 x 842 pt
    page.insert_text((72, 150), "Política de créditos del banco")
    page.insert_text((72, 180), "Este documento describe el proceso de aprobación.")
    # Centrada verticalmente: fuera de las zonas header (top 10%) y footer (bottom 5%),
    # y con resolución natural 120x120 px (sobre el mínimo de 80x80 del filtro)
    page.insert_image(fitz.Rect(200, 350, 320, 470), stream=make_png())
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def make_docx() -> bytes:
    doc = DocxDocument()
    doc.add_heading("Manual de procedimientos", level=1)
    doc.add_paragraph("Primer párrafo del manual interno.")
    doc.add_paragraph("Segundo párrafo con más detalle.")
    doc.add_picture(BytesIO(make_png()), width=DocxInches(2))
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco
    box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(1), PptxInches(6), PptxInches(1)
    )
    box.text_frame.text = "Resultados del trimestre"
    slide.shapes.add_picture(
        BytesIO(make_png()), PptxInches(1), PptxInches(3), width=PptxInches(2)
    )
    slide.notes_slide.notes_text_frame.text = "Recordar mencionar la mora"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_xlsx() -> bytes:
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Cartera"
    ws1.append(["Cliente", "Saldo"])
    ws1.append(["Acme SRL", 15000])
    ws2 = wb.create_sheet("Mora")
    ws2.append(["Cliente", "Días"])
    ws2.append(["Beta SA", 45])
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_image(fmt: str) -> bytes:
    buf = BytesIO()
    Image.new("RGB", (100, 60), (10, 120, 200)).save(buf, format=fmt)
    return buf.getvalue()


def make_rotated_jpeg() -> bytes:
    """JPEG 100x60 con EXIF orientation=6 (rotar 90° CW al mostrar)."""
    img = Image.new("RGB", (100, 60), (10, 120, 200))
    exif = Image.Exif()
    exif[0x0112] = 6
    buf = BytesIO()
    img.save(buf, format="JPEG", exif=exif)
    return buf.getvalue()


# ── PDF ───────────────────────────────────────────────────────────────────


def test_pdf_text_and_image():
    result = pdf_parser.parse(make_pdf())

    all_text = " ".join(b.text for b in result.text_blocks)
    assert "Política de créditos" in all_text
    assert all(b.page_number == 1 for b in result.text_blocks)

    assert len(result.image_blocks) == 1
    assert result.image_blocks[0].data.startswith(PNG_MAGIC)
    assert result.image_blocks[0].page_number == 1


# ── DOCX ──────────────────────────────────────────────────────────────────


def test_docx_text_heading_and_image():
    result = docx_parser.parse(make_docx())

    all_text = "\n".join(b.text for b in result.text_blocks)
    assert "## Manual de procedimientos" in all_text
    assert "Primer párrafo" in all_text
    assert "Segundo párrafo" in all_text

    assert len(result.image_blocks) == 1
    assert len(result.image_blocks[0].data) > 0


# ── PPTX ──────────────────────────────────────────────────────────────────


def test_pptx_text_image_and_notes():
    result = pptx_parser.parse(make_pptx())

    assert len(result.text_blocks) == 1
    block = result.text_blocks[0]
    assert block.page_number == 1
    assert "[Slide 1]" in block.text
    assert "Resultados del trimestre" in block.text
    assert "[Notas del presentador]: Recordar mencionar la mora" in block.text

    assert len(result.image_blocks) == 1
    assert result.image_blocks[0].page_number == 1


# ── XLSX ──────────────────────────────────────────────────────────────────


def test_xlsx_one_block_per_sheet():
    result = xlsx_parser.parse(make_xlsx())

    assert len(result.text_blocks) == 2
    assert result.text_blocks[0].text.startswith("[Hoja: Cartera]")
    assert "Acme SRL\t15000" in result.text_blocks[0].text
    assert result.text_blocks[1].text.startswith("[Hoja: Mora]")
    assert "Beta SA\t45" in result.text_blocks[1].text
    assert [b.page_number for b in result.text_blocks] == [1, 2]


# ── Imágenes sueltas ──────────────────────────────────────────────────────


@pytest.mark.parametrize("fmt", ["PNG", "JPEG", "WEBP"])
def test_image_passthrough(fmt):
    data = make_image(fmt)
    result = image_parser.parse(data)

    assert result.text_blocks == []
    assert len(result.image_blocks) == 1
    # Sin EXIF de rotación los bytes pasan intactos (sin re-encodar)
    assert result.image_blocks[0].data == data


def test_image_exif_rotation_normalized():
    result = image_parser.parse(make_rotated_jpeg())

    out = result.image_blocks[0].data
    assert out.startswith(PNG_MAGIC)
    img = Image.open(BytesIO(out))
    assert img.size == (60, 100)  # dimensiones intercambiadas tras rotar 90°


def test_image_garbage_bytes_pass_through_without_error():
    garbage = b"\x00\x01esto no es una imagen"
    result = image_parser.parse(garbage)
    assert result.image_blocks[0].data == garbage


# ── TXT / MD vía parse_file ───────────────────────────────────────────────


def test_txt_utf8():
    result = parse_file(".txt", "hola señores\n".encode("utf-8"))
    assert result.text_blocks[0].text == "hola señores\n"
    assert result.image_blocks == []


def test_txt_utf8_bom_stripped():
    result = parse_file(".txt", b"\xef\xbb\xbfhola")
    assert result.text_blocks[0].text == "hola"
    assert "﻿" not in result.text_blocks[0].text


def test_md_invalid_bytes_replaced():
    result = parse_file(".md", b"# titulo\n\xff\xfe basura")
    assert result.text_blocks[0].text.startswith("# titulo")
    assert "�" in result.text_blocks[0].text


# ── Dispatch parse_file ───────────────────────────────────────────────────


def test_parse_file_rejects_unsupported_extension():
    with pytest.raises(ValueError, match="Unsupported file extension"):
        parse_file(".zip", b"PK\x03\x04")


def test_every_accepted_extension_has_a_parser():
    fixtures = {
        ".pdf": make_pdf(),
        ".docx": make_docx(),
        ".pptx": make_pptx(),
        ".xlsx": make_xlsx(),
        ".txt": b"texto plano",
        ".md": b"# markdown",
        ".jpg": make_image("JPEG"),
        ".jpeg": make_image("JPEG"),
        ".png": make_image("PNG"),
        ".webp": make_image("WEBP"),
    }
    assert set(fixtures) == ACCEPTED_EXTENSIONS, (
        "Extensión agregada/quitada de ACCEPTED_EXTENSIONS sin actualizar este test"
    )
    for ext, data in fixtures.items():
        result = parse_file(ext, data)
        assert result.text_blocks or result.image_blocks, f"{ext}: parser devolvió vacío"


# ── Archivos corruptos: el parser debe lanzar para que el pipeline marque error ──


@pytest.mark.parametrize(
    "parser",
    [pdf_parser, docx_parser, pptx_parser, xlsx_parser],
    ids=["pdf", "docx", "pptx", "xlsx"],
)
def test_corrupt_file_raises(parser):
    with pytest.raises(Exception):
        parser.parse(b"\x00\x01\x02garbage that is not a valid document")
